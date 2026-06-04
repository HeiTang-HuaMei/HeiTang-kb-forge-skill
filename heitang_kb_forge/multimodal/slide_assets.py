from pathlib import Path

from heitang_kb_forge.multimodal.assets import make_asset
from heitang_kb_forge.schemas.multimodal_schema import MultimodalAsset


def make_slide_assets(path: Path) -> tuple[list[MultimodalAsset], list[dict]]:
    if path.suffix.lower() == ".ppt":
        return [_fallback_slide(path, "Binary PPT parsing is not supported in v1.6.")], []
    try:
        from pptx import Presentation

        presentation = Presentation(path)
    except Exception as exc:
        return [_fallback_slide(path, f"PPTX parsing failed: {exc}")], []

    assets: list[MultimodalAsset] = []
    chunks: list[dict] = []
    for index, slide in enumerate(presentation.slides, start=1):
        texts = []
        title = ""
        for shape in slide.shapes:
            text = getattr(shape, "text", "")
            if text:
                cleaned = text.strip()
                if cleaned:
                    texts.append(cleaned)
                    if not title:
                        title = cleaned.splitlines()[0]
        extracted = "\n".join(texts)
        asset = make_asset(
            path,
            "slide",
            slide_number=index,
            extracted_text=extracted,
            structure={"title": title, "body": extracted, "notes": "", "embedded_assets": []},
            confidence="medium" if extracted else "low",
            extraction_method="parser" if extracted else "fallback",
            review_required=not bool(extracted),
        )
        assets.append(asset)
        if extracted:
            chunks.append(
                {
                    "chunk_id": f"slide_chunk_{asset.asset_id[6:]}",
                    "source_path": str(path).replace("\\", "/"),
                    "source_type": path.suffix.lower().lstrip("."),
                    "title": title or f"Slide {index}",
                    "text": extracted,
                    "order": index - 1,
                    "char_count": len(extracted),
                    "asset_refs": [asset.asset_id],
                    "slide_number": index,
                }
            )
    if not assets:
        return [_fallback_slide(path, "No slides found.")], []
    return assets, chunks


def _fallback_slide(path: Path, message: str) -> MultimodalAsset:
    return make_asset(
        path,
        "slide",
        description=message,
        confidence="low",
        extraction_method="fallback",
        review_required=True,
    )
