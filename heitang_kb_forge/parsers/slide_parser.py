from pathlib import Path


def parse_slide(path: Path) -> str:
    if path.suffix.lower() == ".ppt":
        return ""
    try:
        from pptx import Presentation

        presentation = Presentation(path)
    except Exception:
        return ""
    slide_texts: list[str] = []
    for index, slide in enumerate(presentation.slides, start=1):
        parts = []
        for shape in slide.shapes:
            text = getattr(shape, "text", "")
            if text and text.strip():
                parts.append(text.strip())
        if parts:
            slide_texts.append(f"Slide {index}.\n" + "\n".join(parts))
    return "\n\n".join(slide_texts)
